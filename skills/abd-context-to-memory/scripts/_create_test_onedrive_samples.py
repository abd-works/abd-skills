"""One-off: create sample Word/PowerPoint test trees for context-to-memory testing."""
from __future__ import annotations

import sys
from pathlib import Path


def main() -> int:
    try:
        from docx import Document as DocxDocument
        from pptx import Presentation
    except ImportError:
        print("pip install python-docx python-pptx", file=sys.stderr)
        return 1

    base = Path(r"C:\Users\thoma\OneDrive - Agile by Design\Shared Documents")
    trees = [
        (base / "test-context-to-memory", "Alpha"),
        (base / "test-context-to-memory 2", "Beta"),
    ]

    for root, label in trees:
        root.mkdir(parents=True, exist_ok=True)
        sub_a = root / "notes"
        sub_b = root / "decks"
        sub_a.mkdir(exist_ok=True)
        sub_b.mkdir(exist_ok=True)

        # Root-level short doc
        d = DocxDocument()
        d.add_heading(f"{label} root brief", 0)
        d.add_paragraph(
            f"This is the root document for {root.name}. "
            "Agile by Design context-to-memory pipeline test."
        )
        d.save(str(root / f"{label}_overview.docx"))

        # Subfolder notes
        d2 = DocxDocument()
        d2.add_heading(f"{label} notes", 0)
        d2.add_paragraph("Key point: semantic search should find this phrase about delivery flow.")
        d2.save(str(sub_a / f"{label}_notes.docx"))

        # Subfolder deck (title slide only — avoids layout placeholder quirks)
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = f"{label} deck — RAG test; delivery flow keywords here."
        prs.save(str(sub_b / f"{label}_deck.pptx"))

    rag = base / "test-rag"
    rag.mkdir(parents=True, exist_ok=True)
    (rag / "README.txt").write_text(
        "Aggregate FAISS index (index.faiss, embeddings.npy, metadata.json) "
        "is written here when CONTENT_MEMORY_RAG_PATH points to this folder.\n",
        encoding="utf-8",
    )

    print("Created:", *[str(t[0]) for t in trees], str(rag), sep="\n  ")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
