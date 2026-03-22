"""Create short Word/PowerPoint samples under two OneDrive test roots (ACE context-to-memory E2E)."""
from __future__ import annotations

import sys
from pathlib import Path


def main() -> int:
    try:
        from docx import Document as DocxDocument
        from pptx import Presentation
    except ImportError:
        print("pip install python-docx python-pptx", file=sys.stderr)
        return 1

    base = Path(r"C:\Users\thoma\OneDrive - Agile by Design\Shared Documents")
    trees: list[tuple[Path, str]] = [
        (base / "test-context-to-memory", "Alpha"),
        (base / "test-context-to-memory 2", "Beta"),
    ]

    for root, label in trees:
        root.mkdir(parents=True, exist_ok=True)
        # Root-level docs (2–3)
        d = DocxDocument()
        d.add_heading(f"{label} — root overview", 0)
        d.add_paragraph(
            f"ACE RAG pipeline smoke test. Topic {label}: semantic search should surface "
            "delivery flow, retrospectives, and onboarding checklist phrases."
        )
        d.save(str(root / f"{label}_root_overview.docx"))

        d_one = DocxDocument()
        d_one.add_heading(f"{label} quick facts", 0)
        d_one.add_paragraph(
            "Keywords: agile delivery flow, stakeholder sync, incremental value."
        )
        d_one.save(str(root / f"{label}_quick_facts.docx"))

        prs_root = Presentation()
        slide = prs_root.slides.add_slide(prs_root.slide_layouts[0])
        slide.shapes.title.text = (
            f"{label} — root slide. RAG keywords: onboarding checklist, team charter."
        )
        prs_root.save(str(root / f"{label}_root_deck.pptx"))

        # Subfolders with more short files
        notes = root / "notes"
        decks = root / "decks"
        policies = root / "policies"
        for sub in (notes, decks, policies):
            sub.mkdir(exist_ok=True)

        n1 = DocxDocument()
        n1.add_heading(f"{label} meeting notes", 0)
        n1.add_paragraph(
            "Retro highlights: we improved delivery flow and shortened feedback loops."
        )
        n1.save(str(notes / f"{label}_meeting_notes.docx"))

        n2 = DocxDocument()
        n2.add_heading(f"{label} scratch", 0)
        n2.add_paragraph("Parking lot: explore RAG quality and chunk overlap next sprint.")
        n2.save(str(notes / f"{label}_scratch.docx"))

        prs_d = Presentation()
        s0 = prs_d.slides.add_slide(prs_d.slide_layouts[0])
        s0.shapes.title.text = (
            f"{label} delivery deck — delivery flow and customer outcomes (RAG test)."
        )
        prs_d.save(str(decks / f"{label}_delivery_deck.pptx"))

        pol = DocxDocument()
        pol.add_heading(f"{label} policy snippet", 0)
        pol.add_paragraph(
            "Onboarding checklist: accounts, repo access, runbook link, buddy assignment."
        )
        pol.save(str(policies / f"{label}_policy_snippet.docx"))

    rag = base / "ace-rag-db"
    rag.mkdir(parents=True, exist_ok=True)
    (rag / "README.txt").write_text(
        "Aggregate FAISS index (index.faiss, embeddings.npy, metadata.json) "
        "when CONTENT_MEMORY_RAG_PATH points here and embed_and_index runs from the hub.\n",
        encoding="utf-8",
    )

    print("Created sample trees and ace-rag-db placeholder:")
    for root, _ in trees:
        print(" ", root)
    print(" ", rag)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
